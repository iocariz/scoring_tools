"""PowerPoint companion for score_rebin.py — generic, data-driven deck.

`build_deck(args, schemes, mono, ginis, wlabel)` is called by score_rebin.py when `--deck` is
passed. Everything is read from the computed result objects, so it works for any population/score.
Requires python-pptx (already a project dependency).
"""

from __future__ import annotations

import os

import numpy as np
from loguru import logger
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(0x1F, 0x35, 0x5E)
GREEN = RGBColor(0x2C, 0xA0, 0x2C)
GREY = RGBColor(0x66, 0x66, 0x66)


def _fmt(e: float) -> str:
    return "-inf" if e == -np.inf else "+inf" if e == np.inf else f"{e:.7f}"


def _pct(v) -> str:
    return "n/a" if v is None or (isinstance(v, float) and np.isnan(v)) else f"{v:.2f}%"


def _reversals(b2) -> list[int]:
    a = np.asarray(b2, dtype=float)
    return [i + 1 for i in range(1, len(a)) if not np.isnan(a[i]) and not np.isnan(a[i - 1]) and a[i] > a[i - 1] + 1e-9]


def _textbox(s, left, top, w, h):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _heading(s, text, sub=None):
    tf = _textbox(s, 0.5, 0.3, 12.3, 1.1)
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size, r.font.bold, r.font.color.rgb = Pt(27), True, NAVY
    if sub:
        p = tf.add_paragraph()
        r2 = p.add_run()
        r2.text = sub
        r2.font.size, r2.font.color.rgb = Pt(13), GREY


def _bullets(s, items, top=1.7, size=15):
    tf = _textbox(s, 0.7, top, 12.0, 5.2)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = "• " + it
        r.font.size = Pt(size)
        p.space_after = Pt(6)


def _table(s, data, left, top, w, h, col_w=None, fs=11):
    gf = s.shapes.add_table(len(data), len(data[0]), Inches(left), Inches(top), Inches(w), Inches(h))
    t = gf.table
    if col_w:
        for j, cw in enumerate(col_w):
            t.columns[j].width = Inches(cw)
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            c = t.cell(i, j)
            c.text = str(val)
            para = c.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            rn = para.runs[0]
            rn.font.size = Pt(fs)
            c.fill.solid()
            if i == 0:
                rn.font.bold = True
                rn.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                c.fill.fore_color.rgb = NAVY
            else:
                c.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if i % 2 else RGBColor(0xEE, 0xF1, 0xF7)


def build_deck(args, schemes, mono, ginis, wlabel) -> str:
    primary = args.primary_k
    sp = schemes[primary]
    edges = sp.attrs["edges"]
    arr = "[" + ", ".join(_fmt(e) for e in edges) + "]"

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    # ---- 1. title + recommendation ----
    s = prs.slides.add_slide(blank)
    tf = _textbox(s, 0.7, 1.4, 12.0, 1.8)
    r = tf.paragraphs[0].add_run()
    r.text = args.title
    r.font.size, r.font.bold, r.font.color.rgb = Pt(38), True, NAVY
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = f"Score `{args.score_col}`  ·  weight {wlabel}  ·  {len(sp)} groups  ·  {args.date_tag}"
    r.font.size, r.font.color.rgb = Pt(15), GREY
    gini_map = dict(ginis["schemes"]) if ginis else {}
    prim_gini = gini_map.get(f"K={primary} balanced")
    box = _textbox(s, 0.7, 3.7, 12.0, 3.0)
    r = box.paragraphs[0].add_run()
    r.text = f"Recommendation — K={primary}"
    r.font.size, r.font.bold, r.font.color.rgb = Pt(18), True, GREEN
    for line in [
        f"Weight balance {sp['weight_share'].min():.1%}–{sp['weight_share'].max():.1%} "
        f"(std {sp['weight_share'].std():.3f}, target {1 / len(sp):.1%}).",
        (
            f"Gini retained {prim_gini / ginis['continuous'] * 100:.1f}%."
            if prim_gini is not None
            else "Risk (b2_ever_h6) reported per group."
        ),
        f"{args.score_col}_bins = {arr}",
    ]:
        pp = box.add_paragraph()
        rr = pp.add_run()
        rr.text = "✓ " + line
        rr.font.size = Pt(13 if line.startswith(args.score_col) else 14)
        if line.startswith(args.score_col):
            rr.font.name = "Courier New"

    # ---- 2. scheme comparison ----
    s = prs.slides.add_slide(blank)
    _heading(s, "Scheme comparison", "Balanced (merge + split oversized) vs risk-monotone (isotonic).")
    rows = [["Scheme", "K", f"{wlabel} balance (std)", "Max group", "b2 reversals"]]
    for k in args.groups:
        g = schemes[k]
        revs = _reversals(g["b2_ever_h6"])
        rows.append(
            [
                "balanced",
                str(k),
                f"{g['weight_share'].std():.3f}",
                f"{g['weight_share'].max():.1%}",
                str(revs) if revs else "none",
            ]
        )
    if mono is not None:
        rows.append(
            [
                "risk-monotone",
                str(len(mono)),
                f"{mono['weight_share'].std():.3f}",
                f"{mono['weight_share'].max():.1%}",
                "none (strict)",
            ]
        )
    _table(s, rows, 0.7, 2.3, 11.9, 2.4, col_w=[3.0, 1.0, 3.0, 2.4, 2.5], fs=13)

    # ---- 3. recommended K group table + array ----
    s = prs.slides.add_slide(blank)
    _heading(
        s,
        f"Recommended scheme — K={primary}",
        f"{wlabel} band {sp['weight_share'].min():.1%}–{sp['weight_share'].max():.1%}.",
    )
    gdata = [["Grp", "Range", f"{wlabel} share", "b2_ever_h6"]]
    for _, rr in sp.iterrows():
        gdata.append([int(rr["group"]), rr["range"], f"{rr['weight_share']:.1%}", _pct(rr["b2_ever_h6"])])
    _table(s, gdata, 0.6, 1.7, 7.4, 5.3, col_w=[0.7, 3.6, 1.5, 1.6], fs=10)
    tf = _textbox(s, 8.2, 1.8, 4.7, 5.0)
    r = tf.paragraphs[0].add_run()
    r.text = f"{args.score_col}_bins (K={primary})"
    r.font.size, r.font.bold, r.font.color.rgb = Pt(13), True, NAVY
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = arr
    r.font.size, r.font.name = Pt(10), "Courier New"

    # ---- 4. conversion map ----
    s = prs.slides.add_slide(blank)
    _heading(s, "Conversion map: old candidate bin → new group")
    cdata = [["Old bin(s)", "Range", "→ New group"]]
    for _, rr in sp.iterrows():
        fb, lb = int(rr["first_bin"]), int(rr["last_bin"])
        cdata.append([str(fb) if fb == lb else f"{fb}–{lb}", rr["range"], int(rr["group"])])
    _table(s, cdata, 0.7, 1.7, 11.9, 5.2, col_w=[2.4, 5.5, 4.0], fs=11)

    # ---- 5. plot ----
    img = f"{args.output_dir}/{args.prefix}_balance_k{primary}.png"
    if os.path.exists(img):
        s = prs.slides.add_slide(blank)
        _heading(s, f"K={primary} — balance + realized risk (b2)")
        s.shapes.add_picture(img, Inches(1.4), Inches(1.3), height=Inches(5.9))

    # ---- 6. Gini ----
    if ginis:
        s = prs.slides.add_slide(blank)
        _heading(s, "Discriminance (Gini vs target)", f"{ginis['n']:,} with target, bad-rate {ginis['bad_rate']:.2%}.")
        gd = [["Scheme", "Gini", "Retention"], ["continuous score", f"{ginis['continuous']:.4f}", "100.0%"]]
        for name, g in ginis["schemes"]:
            gd.append([name, f"{g:.4f}", f"{g / ginis['continuous'] * 100:.1f}%"])
        _table(s, gd, 0.7, 2.0, 8.0, 4.0, col_w=[4.0, 1.8, 2.0], fs=12)
        _bullets(
            s,
            [
                "Binning is near-lossless when retention stays >98%.",
                "The ceiling is the score itself, not the grouping.",
            ],
            top=6.2,
            size=13,
        )

    out = f"{args.output_dir}/{args.prefix}_deck.pptx"
    prs.save(out)
    logger.info(f"Saved deck {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
    return out
