"""Generate a management **results** presentation (.pptx) from a completed batch run.

Distinct from ``generate_presentation.py`` (which explains the *methodology*): this builds an
executive **results** deck from the consolidated outputs — headline production/risk, per-segment
charts, the swap-in/swap-out bridge, acceptance change, the audit-category risk surfaces (2D
heatmap + 3D snapshot), and the out-of-time validation — for a management audience.

Usage
-----
    uv run python generate_results_presentation.py                       # uses output/
    uv run python generate_results_presentation.py --surface-segment premium
    uv run python generate_results_presentation.py --scenario base --output output/results_deck
    uv run python generate_results_presentation.py --pdf                 # also export PDF (needs LibreOffice)
"""

from __future__ import annotations

import argparse
import subprocess
import sys
from pathlib import Path

import pandas as pd
from loguru import logger
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from plot_risk_surface import _discover_segments, _load_unit, _reporting_supersegments
from src import presentation_charts as pc
from src.risk_surface import aggregate_to_cells, build_risk_surface_figure, classify_cells

NAVY = RGBColor(0x2C, 0x3E, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x34, 0x98, 0xDB)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
RED = RGBColor(0xE7, 0x4C, 0x3C)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xEA, 0xF1, 0xF8)


# --------------------------------------------------------------------------- #
# pptx helpers
# --------------------------------------------------------------------------- #


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _box(slide, left, top, width, height, lines, *, size=16, color=NAVY, bold=False, align=PP_ALIGN.LEFT, anchor=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, line in enumerate([lines] if isinstance(lines, str) else lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def _bullets(slide, left, top, width, height, items, *, size=15):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (txt, col) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = "•  " + txt
        run.font.size = Pt(size)
        run.font.color.rgb = col
        run.font.name = "Calibri"
    return tb


def _title_bar(slide, prs, title, subtitle=""):
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False
    _box(slide, 0.5, 0.08, 12.3, 0.6, title, size=26, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        _box(slide, 0.5, 0.62, 12.3, 0.34, subtitle, size=12, color=RGBColor(0xC9, 0xD6, 0xE5))


def _image(slide, path: Path, *, left=None, top=1.15, width=None, height=None, prs=None):
    kw = {}
    if width is not None:
        kw["width"] = Inches(width)
    if height is not None:
        kw["height"] = Inches(height)
    pic = slide.shapes.add_picture(str(path), Inches(0), Inches(top), **kw)
    if left is None and prs is not None:  # centre horizontally
        pic.left = int((prs.slide_width - pic.width) / 2)
    elif left is not None:
        pic.left = Inches(left)
    return pic


def _kpi_card(slide, left, top, width, label, value, sub="", accent=ACCENT):
    card = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(1.5))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = accent
    card.line.width = Pt(1.25)
    card.shadow.inherit = False
    _box(slide, left + 0.1, top + 0.12, width - 0.2, 0.3, label, size=12, color=GREY)
    _box(slide, left + 0.1, top + 0.42, width - 0.2, 0.6, value, size=24, color=accent, bold=True)
    if sub:
        _box(slide, left + 0.1, top + 1.08, width - 0.2, 0.34, sub, size=11, color=GREY)


# --------------------------------------------------------------------------- #
# Data helpers
# --------------------------------------------------------------------------- #


def _fmt_eur(v: float) -> str:
    return f"€{v / 1e6:,.1f}M"


def _all_surface_units(
    data_dir: Path, segments_config: Path, suffix: str, only: str | None = None, include_supersegments: bool = True
) -> list[tuple[str, pd.DataFrame, list[str]]]:
    """Enumerate (label, classified-cells, variables) for every segment (+ reporting supersegment).

    Reuses the same discovery + supersegment aggregation as ``plot_risk_surface``; ``only`` restricts
    to a single segment.
    """
    seg_dirs = _discover_segments(Path(data_dir), suffix, [only] if only else None)
    loaded: dict[str, tuple] = {}
    units: list[tuple[str, pd.DataFrame, list[str]]] = []
    for seg_dir in seg_dirs:
        u = _load_unit(seg_dir, suffix)
        if u is None:
            continue
        loaded[seg_dir.name] = u
        cells = aggregate_to_cells(classify_cells(u[0], u[1], u[2]), u[2], u[3])
        if not cells.empty:
            units.append((seg_dir.name, cells, u[2]))
    if include_supersegments and not only:
        for name, members in _reporting_supersegments(Path(segments_config), set(loaded)).items():
            us = [loaded[m] for m in members]
            if len({tuple(x[2]) for x in us}) != 1:
                continue
            variables, mult = us[0][2], us[0][3]
            long = pd.concat([classify_cells(s, a, variables) for s, a, _v, _m in us], ignore_index=True)
            cells = aggregate_to_cells(long, variables, mult)
            if not cells.empty:
                units.append((f"{name} (supersegment)", cells, variables))
    return units


# --------------------------------------------------------------------------- #
# Build
# --------------------------------------------------------------------------- #


def build_results_presentation(
    data_dir: str | Path = "output",
    out_dir: str | Path = "output/results_deck",
    scenario: str = "base",
    surface_segment: str | None = None,
    segments_config: str | Path = "segments.toml",
    today: str = "",
) -> Path:
    """Build the management results deck. Returns the saved .pptx path."""
    data_dir, out_dir = Path(data_dir), Path(out_dir)
    img_dir = out_dir / "images"
    img_dir.mkdir(parents=True, exist_ok=True)
    suffix = f"_{scenario}" if scenario else ""

    cons = pd.read_csv(data_dir / "consolidated_risk_production.csv")
    bt_path = data_dir / "backtest" / f"backtest_consolidated{suffix}.csv"
    bt = pd.read_csv(bt_path) if bt_path.exists() else None
    total = pc._total_row(cons)

    # --- charts ---
    charts = {
        "production": pc.chart_production_by_segment,
        "risk": pc.chart_risk_by_segment,
        "swap": pc.chart_swap_waterfall,
        "acceptance": pc.chart_acceptance_by_segment,
    }
    paths = {k: pc.save_chart(fn(cons), img_dir / f"{k}.png") for k, fn in charts.items()}
    if bt is not None and not bt.empty:
        paths["backtest"] = pc.save_chart(pc.chart_backtest_validation(bt), img_dir / "backtest.png")

    # --- risk surfaces for ALL segments (+ reporting supersegments) ---
    surface_slides: list[tuple[str, Path | None, Path | None]] = []
    for label, cells, variables in _all_surface_units(data_dir, segments_config, suffix, only=surface_segment):
        facet = variables[2] if len(variables) > 2 else None
        safe = label.split(" ")[0].replace("/", "_")
        hp = pc.save_chart(
            pc.chart_category_heatmap(cells, variables[:2], facet, f"Risk surface — {label}"),
            img_dir / f"heatmap_{safe}.png",
        )
        fig3d = build_risk_surface_figure(cells, variables[:2], facet, f"Risk surface — {label}")
        sp = pc.export_surface_png(fig3d, img_dir / f"surface3d_{safe}.png")
        surface_slides.append((label, hp, sp))

    # --- deck ---
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    _slide_title(prs, total, scenario, today)
    _slide_exec_summary(prs, total, bt)
    _slide_methodology(prs)
    _slide_chart(
        prs,
        "Production by segment",
        paths["production"],
        "Proposed cutoffs vs the current booked portfolio. Several segments are tightened to reach the risk target.",
    )
    _slide_chart(
        prs,
        "Risk by segment",
        paths["risk"],
        "Annualised b2_ever_h6. The proposal brings each segment to (or below) its target risk appetite.",
    )
    _slide_chart(
        prs,
        "Production bridge — swap-in vs swap-out",
        paths["swap"],
        "Swap-in = newly accepted (previously rejected); swap-out = now rejected (previously booked).",
    )
    _slide_chart(
        prs,
        "Acceptance rate by segment",
        paths["acceptance"],
        "Share of demand accepted under the current vs proposed cutoffs.",
    )
    for label, hp, sp in surface_slides:
        if hp or sp:
            _slide_surfaces(prs, label, hp, sp)
    if "backtest" in paths:
        _slide_chart(
            prs,
            "Out-of-time validation",
            paths["backtest"],
            "Frozen cutoffs replayed on a matured held-out cohort. Flag is noise-aware (OK / INCONCLUSIVE / DRIFT).",
        )
    _slide_next_steps(prs, total)

    out = out_dir / f"Credit_Risk_Results_{scenario}.pptx"
    prs.save(str(out))
    logger.info(f"Wrote {out} ({len(prs.slides)} slides)")
    return out


def _slide_title(prs, total, scenario, today):
    slide = _blank(prs)
    rect = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = NAVY
    rect.line.fill.background()
    rect.shadow.inherit = False
    _box(slide, 0.8, 2.5, 11.7, 1.0, "Credit Risk Cutoff Optimisation", size=40, color=WHITE, bold=True)
    _box(
        slide,
        0.8,
        3.5,
        11.7,
        0.6,
        "Results & Recommendation — for Management",
        size=22,
        color=RGBColor(0x9B, 0xC4, 0xE8),
    )
    sub = f"Scenario: {scenario}" + (f"  ·  {today}" if today else "")
    if total is not None:
        sub += f"  ·  {int(total.get('n_segments', 0))} segments"
    _box(slide, 0.8, 4.4, 11.7, 0.5, sub, size=14, color=RGBColor(0xC9, 0xD6, 0xE5))


def _slide_exec_summary(prs, total, bt):
    slide = _blank(prs)
    _title_bar(slide, prs, "Executive summary", "Current booked portfolio vs proposed cutoffs (base scenario)")
    if total is not None:
        ap, op = total["actual_production"], total["optimum_production"]
        ar, orr = total["actual_risk_pct"], total["optimum_risk_pct"]
        dpct = total.get("production_delta_pct", (op - ap) / ap * 100 if ap else 0)
        _kpi_card(slide, 0.5, 1.4, 3.0, "Production (current)", _fmt_eur(ap), accent=GREY)
        _kpi_card(
            slide,
            3.7,
            1.4,
            3.0,
            "Production (proposed)",
            _fmt_eur(op),
            f"{dpct:+.1f}% vs current",
            accent=GREEN if dpct >= 0 else RED,
        )
        _kpi_card(slide, 6.9, 1.4, 3.0, "Risk (current → proposed)", f"{ar:.2f}% → {orr:.2f}%", accent=RED)
        _kpi_card(slide, 10.1, 1.4, 2.7, "Net risk change", f"{orr - ar:+.2f} pp", accent=ACCENT)
    bullets = _recommendation_bullets(total, bt)
    _box(slide, 0.5, 3.3, 12.3, 0.4, "Recommendation & key points", size=16, color=NAVY, bold=True)
    _bullets(slide, 0.6, 3.8, 12.2, 3.3, bullets, size=15)


def _recommendation_bullets(total, bt):
    out = []
    if total is not None:
        ar, orr = total["actual_risk_pct"], total["optimum_risk_pct"]
        dpct = total.get("production_delta_pct", 0)
        verb = "lifts" if dpct >= 0 else "trades"
        out.append(
            (
                f"Proposed cutoffs bring the portfolio risk to {orr:.2f}% (from {ar:.2f}%) and {verb} "
                f"production by {dpct:+.1f}% — the optimal risk/production point at the agreed appetite.",
                NAVY,
            )
        )
        out.append(
            (
                "Gains come from swap-in (previously score-rejected, now accepted); reductions come from "
                "swapping out the riskiest currently-booked cells.",
                GREY,
            )
        )
    if bt is not None and not bt.empty and "drift_flag" in bt.columns:
        vc = bt["drift_flag"].value_counts().to_dict()
        out.append(
            (
                f"Out-of-time validation: {vc.get('OK', 0)} OK, {vc.get('INCONCLUSIVE', 0)} inconclusive, "
                f"{vc.get('DRIFT', 0)} drift across {len(bt)} backtested segments (noise-aware).",
                GREY,
            )
        )
    out.append(
        (
            "Every figure carries bootstrap confidence intervals; cutoffs are reproducible and registered "
            "(champion policy per segment).",
            GREY,
        )
    )
    out.append(("Recommendation: adopt the proposed cutoffs, monitor the swap-in segments as cohorts mature.", ACCENT))
    return out


def _slide_methodology(prs):
    slide = _blank(prs)
    _title_bar(slide, prs, "How these numbers are produced", "Methodology in brief")
    items = [
        (
            "Populations: booked loans (observed outcomes), score-rejected 'repesca' (outcomes inferred), and "
            "policy-rejected (excluded — their rejection is unrelated to the credit score).",
            NAVY,
        ),
        (
            "Risk metric: b2_ever_h6 — the annualised, exposure-weighted 6-month delinquency rate. Repesca risk is "
            "predicted by a trained model, then corrected for selection bias (stress factor + reject inference).",
            GREY,
        ),
        (
            "Cutoffs: a Mixed-Integer Linear Program maximises booked production within a risk budget, under "
            "monotonicity (a 'staircase' — never accept a riskier cell while rejecting a safer one).",
            GREY,
        ),
        (
            "A menu, not a point: a Pareto frontier of risk vs production is built; three operating points "
            "(pessimistic / base / optimistic) bracket the risk appetite.",
            GREY,
        ),
        (
            "Validation & governance: out-of-time backtest on matured cohorts, bootstrap confidence intervals, "
            "PSI/CSI stability, and a reproducibility + champion-policy registry.",
            GREY,
        ),
    ]
    _bullets(slide, 0.7, 1.4, 12.0, 5.2, items, size=15)
    _box(
        slide,
        0.6,
        6.7,
        12.1,
        0.5,
        "Booked → bin scores → infer repesca risk → MILP cutoffs → Pareto scenarios → out-of-time validation",
        size=13,
        color=ACCENT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def _slide_chart(prs, title, img_path, caption=""):
    slide = _blank(prs)
    _title_bar(slide, prs, title)
    _image(slide, img_path, top=1.25, width=11.6, prs=prs)
    if caption:
        _box(slide, 0.6, 6.95, 12.1, 0.5, caption, size=12, color=GREY, align=PP_ALIGN.CENTER)


def _slide_surfaces(prs, seg_name, heatmap_path, surf3d_path):
    slide = _blank(prs)
    _title_bar(
        slide,
        prs,
        "Risk surfaces by audit category",
        f"{seg_name}: b2_ever_h6 over the two score bins, coloured by keep / swap-in / swap-out / rejected",
    )
    if heatmap_path and surf3d_path:
        _image(slide, heatmap_path, left=0.4, top=1.4, width=6.3)
        _image(slide, surf3d_path, left=6.9, top=1.4, width=6.1)
    elif heatmap_path:
        _image(slide, heatmap_path, top=1.4, width=8.5, prs=prs)
    elif surf3d_path:
        _image(slide, surf3d_path, top=1.4, width=8.5, prs=prs)
    _box(
        slide,
        0.6,
        6.85,
        12.1,
        0.5,
        "Category = (booked-before × accepted-now). 2D heatmap (left) annotates risk %; 3D surface (right) shows the height.",
        size=12,
        color=GREY,
        align=PP_ALIGN.CENTER,
    )


def _slide_next_steps(prs, total):
    slide = _blank(prs)
    _title_bar(slide, prs, "Recommendation & next steps")
    items = [
        ("Adopt the proposed base-scenario cutoffs as the new policy (one champion per segment).", NAVY),
        ("Register the policy and run the champion/challenger comparison each refresh to catch drift early.", GREY),
        ("Monitor swap-in segments out-of-time as their cohorts mature (the validation flag turns OK/DRIFT).", GREY),
        ("Re-validate on any new data snapshot (reproducibility check fails loudly if inputs changed).", GREY),
        ("Optional: review the pessimistic / optimistic scenarios for the risk-appetite trade-off band.", ACCENT),
    ]
    _bullets(slide, 0.7, 1.5, 12.0, 4.5, items, size=17)


def convert_to_pdf(pptx_path: Path) -> Path | None:
    try:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir", str(pptx_path.parent), str(pptx_path)],
            check=True,
            capture_output=True,
            timeout=120,
        )
        pdf = pptx_path.with_suffix(".pdf")
        return pdf if pdf.exists() else None
    except Exception as e:
        logger.warning(f"PDF conversion skipped ({type(e).__name__}); LibreOffice (soffice) required.")
        return None


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Generate the management results presentation (.pptx).",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    parser.add_argument("--data-dir", default="output", help="Dir with consolidated outputs + segment runs.")
    parser.add_argument("--output", "-o", default="output/results_deck", help="Output dir for the deck + images.")
    parser.add_argument("--scenario", default="base", help="Scenario suffix (default: base).")
    parser.add_argument(
        "--surface-segment", default=None, help="Restrict risk-surface slides to one segment (default: all)."
    )
    parser.add_argument("--segments-config", default="segments.toml", help="For reporting-supersegment surfaces.")
    parser.add_argument("--today", default="", help="Date string for the title slide.")
    parser.add_argument("--pdf", action="store_true", help="Also export a PDF (requires LibreOffice).")
    args = parser.parse_args(argv if argv is not None else sys.argv[1:])

    cons = Path(args.data_dir) / "consolidated_risk_production.csv"
    if not cons.exists():
        logger.error(f"Consolidated results not found: {cons} — run run_batch.py first.")
        return 1

    out = build_results_presentation(
        data_dir=args.data_dir,
        out_dir=args.output,
        scenario=args.scenario,
        surface_segment=args.surface_segment,
        segments_config=args.segments_config,
        today=args.today,
    )
    if args.pdf:
        pdf = convert_to_pdf(out)
        if pdf:
            logger.info(f"Wrote {pdf}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
