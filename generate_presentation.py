"""Generate the Credit Risk Scoring Pipeline Presentation as .pptx and .pdf.

Usage:
    uv run python generate_presentation.py          # generates .pptx
    uv run python generate_presentation.py --pdf     # also converts to .pdf (requires LibreOffice)
"""

from __future__ import annotations

import argparse
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Brand palette
# ---------------------------------------------------------------------------
DARK_BLUE = RGBColor(0x2E, 0x52, 0x8F)  # header bg
MID_BLUE = RGBColor(0x3A, 0x6E, 0xB5)   # accent
LIGHT_BLUE = RGBColor(0x4A, 0x90, 0xD9)  # lighter accent
GREEN = RGBColor(0x2E, 0xCC, 0x71)       # accent bar
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x2C, 0x2C, 0x2C)
GREY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GREY = RGBColor(0xF0, 0xF0, 0xF0)
TABLE_HEADER_BG = RGBColor(0x2E, 0x3B, 0x55)
TABLE_ALT_BG = RGBColor(0xF5, 0xF7, 0xFA)
ACCEPT_GREEN = RGBColor(0xD4, 0xED, 0xDA)
REJECT_RED = RGBColor(0xF8, 0xD7, 0xDA)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_size=14,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def _add_para(tf, text, font_size=14, color=BLACK, bold=False, space_before=Pt(4),
              alignment=PP_ALIGN.LEFT, font_name="Calibri", bullet=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    if bullet:
        p.level = 0
    return p


def _draw_header(slide, title, subtitle=""):
    """Draw the blue header banner + green accent bar."""
    # Blue header bg
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.7)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title
    _add_textbox(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.7),
                 title, font_size=36, color=WHITE, bold=True)

    # Subtitle
    if subtitle:
        _add_textbox(slide, Inches(0.8), Inches(1.05), Inches(11), Inches(0.5),
                     subtitle, font_size=16, color=RGBColor(0xBB, 0xCC, 0xDD))

    # Green accent bar
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(1.7), SLIDE_WIDTH, Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()


def _draw_footer(slide, slide_num, total_slides):
    _add_textbox(slide, Inches(0.5), Inches(7.05), Inches(6), Inches(0.3),
                 f"Credit Risk Scoring & Portfolio Optimization  |  Slide {slide_num}/{total_slides}",
                 font_size=8, color=GREY)
    _add_textbox(slide, Inches(10.5), Inches(7.05), Inches(2.5), Inches(0.3),
                 "Confidential", font_size=8, color=GREY, alignment=PP_ALIGN.RIGHT)


def _add_table(slide, left, top, width, rows_data, col_widths=None,
               font_size=11, header_font_size=12):
    """Add a styled table. rows_data[0] is the header row."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.35 * n_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row_data in enumerate(rows_data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(header_font_size if r == 0 else font_size)
                paragraph.font.color.rgb = WHITE if r == 0 else BLACK
                paragraph.font.bold = r == 0
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG

    return table_shape


def _section_heading(tf, text, font_size=18):
    _add_para(tf, text, font_size=font_size, color=DARK_BLUE, bold=True, space_before=Pt(12))


def _bullet(tf, text, font_size=13, bold_prefix="", color=BLACK):
    if bold_prefix:
        p = tf.add_paragraph()
        run1 = p.add_run()
        run1.text = bold_prefix + " "
        run1.font.size = Pt(font_size)
        run1.font.bold = True
        run1.font.color.rgb = color
        run1.font.name = "Calibri"
        run2 = p.add_run()
        run2.text = text
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = color
        run2.font.name = "Calibri"
        p.space_before = Pt(3)
    else:
        p = _add_para(tf, f"  {text}", font_size=font_size, color=color, space_before=Pt(3))
    return p


def _code_box(slide, left, top, width, height, text, font_size=11):
    """Monospaced code block with light background."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GREY
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = "Consolas"
        p.font.color.rgb = BLACK
        p.space_before = Pt(1)
    return tf


def _flow_box(slide, left, top, width, height, text, bg_color=LIGHT_BLUE):
    """Colored flow/phase box."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(2)
    return shape


def _arrow(slide, left, top):
    """Small right arrow."""
    _add_textbox(slide, left, top, Inches(0.3), Inches(0.35),
                 "\u2192", font_size=18, color=GREY, alignment=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

TOTAL_SLIDES = 18


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, DARK_BLUE)

    _add_textbox(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.5),
                 "Credit Risk Scoring\n& Portfolio Optimization",
                 font_size=48, color=WHITE, bold=True)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(3.0), SLIDE_WIDTH, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()

    _add_textbox(slide, Inches(1), Inches(3.4), Inches(11), Inches(0.5),
                 "Technical Methodology & Architecture",
                 font_size=22, color=RGBColor(0xCC, 0xDD, 0xEE))

    _add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(1.0),
                 "End-to-end pipeline for multi-segment credit decisioning\n"
                 "with MILP optimization, reject inference, and automated reporting",
                 font_size=16, color=RGBColor(0xAA, 0xBB, 0xCC))

    _add_textbox(slide, Inches(1), Inches(5.8), Inches(6), Inches(0.4),
                 "March 2026  |  Confidential",
                 font_size=14, color=RGBColor(0x99, 0xAA, 0xBB))

    _draw_footer(slide, 1, TOTAL_SLIDES)


def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_header(slide, "Agenda", "Presentation outline")

    items = [
        "Pipeline Architecture & End-to-End Flow",
        "Data Populations & Risk Metric",
        "Preprocessing & Binning Methodology",
        "Risk Adjustment Pipeline (Stress, RI, Tasa Fin)",
        "Model Inference (Hurdle, Tweedie, XGBoost, LightGBM)",
        "MILP Optimization & Pareto Frontier",
        "Reject Inference & Selection Bias Correction",
        "Scenario Analysis & Bootstrap CI",
        "Stability (PSI/CSI) & MR Validation",
        "Sensitivity, Marginal Impact & Fixed-Cell Constraints",
        "Batch Processing & Global Allocation",
        "Trend Analysis & SPC Anomaly Detection",
        "Consolidated Reporting",
        "Recent Improvements & Summary",
    ]

    tf = _add_textbox(slide, Inches(1.5), Inches(2.1), Inches(10), Inches(5),
                      "", font_size=15)
    for i, item in enumerate(items):
        color = LIGHT_BLUE if (i + 1) % 2 == 0 else DARK_BLUE
        p = _add_para(tf, f"  {i + 1}.   {item}", font_size=15, color=BLACK,
                       space_before=Pt(6))

    _draw_footer(slide, 2, TOTAL_SLIDES)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_header(slide, "Pipeline Architecture", "End-to-end flow from raw data to decisions")

    # Phase boxes - row 1
    phases_r1 = [
        "1. Config\n(Pydantic)",
        "2. Data Load\n(SAS files)",
        "3. Preprocess\n(DQ + Bins)",
        "4. Inference\n(Models + CV)",
    ]
    y1 = Inches(2.2)
    for i, label in enumerate(phases_r1):
        x = Inches(0.8 + i * 3.1)
        _flow_box(slide, x, y1, Inches(2.5), Inches(0.7), label, bg_color=LIGHT_BLUE)
        if i < 3:
            _arrow(slide, Inches(0.8 + (i + 1) * 3.1 - 0.4), y1 + Inches(0.15))

    # Arrow down
    _add_textbox(slide, Inches(11.3), Inches(2.95), Inches(0.5), Inches(0.4),
                 "\u2193", font_size=20, color=GREY, alignment=PP_ALIGN.CENTER)

    # Phase boxes - row 2
    phases_r2 = [
        "5. Optimization\n(MILP / Pareto)",
        "6. Scenarios\n(Pessim/Base/Opt)",
        "7. Stability\n(PSI + MR)",
        "8. Reporting\n(HTML)",
    ]
    y2 = Inches(3.5)
    for i, label in enumerate(phases_r2):
        x = Inches(0.8 + i * 3.1)
        _flow_box(slide, x, y2, Inches(2.5), Inches(0.7), label, bg_color=MID_BLUE)
        if i < 3:
            _arrow(slide, Inches(0.8 + (i + 1) * 3.1 - 0.4), y2 + Inches(0.15))

    # Optional extensions
    tf = _add_textbox(slide, Inches(0.8), Inches(4.5), Inches(6), Inches(2.5), "", font_size=13)
    _section_heading(tf, "Optional Pipeline Extensions", font_size=16)
    _bullet(tf, "Sensitivity analysis: risk perturbation across grid", bold_prefix="6b.")
    _bullet(tf, "RI parameter optimization: grid search / Optuna TPE", bold_prefix="6c.")
    _bullet(tf, "Re-optimization with tuned RI parameters", bold_prefix="6d.")
    _bullet(tf, "Trend analysis: monthly SPC anomaly detection", bold_prefix="7b.")

    # Design principles
    tf2 = _add_textbox(slide, Inches(7), Inches(4.5), Inches(6), Inches(2.5), "", font_size=13)
    _section_heading(tf2, "Key Design Principles", font_size=16)
    _bullet(tf2, "Centralized OutputPaths dataclass for all file I/O", bold_prefix="Path Mgmt:")
    _bullet(tf2, "Non-blocking optional steps (trends, reporting, sensitivity)", bold_prefix="Isolation:")
    _bullet(tf2, "PreprocessingResult dataclass for phase outputs", bold_prefix="Typed Returns:")
    _bullet(tf2, "Training-only mode, pre-trained model loading, skip DQ", bold_prefix="Early Exit:")

    _draw_footer(slide, 3, TOTAL_SLIDES)


def slide_populations_risk(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_header(slide, "Data Populations & Risk Metric", "Three populations, one risk measure")

    tf = _add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.5), "", font_size=13)
    _section_heading(tf, "Application Populations", font_size=18)

    _add_table(slide, Inches(0.8), Inches(2.6), Inches(11.5), [
        ["Population", "Definition", "Role in Optimization"],
        ["Booked", "Approved and disbursed applications", "Observed outcomes (risk + production)"],
        ["Score-rejected (repesca)", "Rejected by credit score cutoff (09-score)", "Candidates for cutoff changes; outcomes inferred"],
        ["Policy-rejected", "Rejected for non-score reasons (08-other)", "Excluded -- rejection unrelated to creditworthiness"],
    ])

    tf2 = _add_textbox(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(0.5), "", font_size=13)
    _section_heading(tf2, "Risk Metric: b2_ever_h6", font_size=18)

    _code_box(slide, Inches(0.8), Inches(5.1), Inches(7), Inches(0.8),
              "b2_ever_h6 = multiplier x todu_30ever_h6 / todu_amt_pile_h6\n"
              "           (default multiplier = 7 for H6, 4 for H3)", font_size=12)

    tf3 = _add_textbox(slide, Inches(8.5), Inches(4.5), Inches(4.5), Inches(2.5), "", font_size=13)
    _section_heading(tf3, "Annualization", font_size=16)
    _add_para(tf3, "All indicators scaled by:", font_size=12, space_before=Pt(6))
    _add_para(tf3, "annual_coef = 12 / n_months", font_size=12, bold=True, space_before=Pt(4))
    _add_para(tf3, "Normalizes to 12-month equivalent across different observation periods.",
              font_size=11, color=GREY, space_before=Pt(4))

    _draw_footer(slide, 4, TOTAL_SLIDES)


def slide_preprocessing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_header(slide, "Preprocessing & Binning", "Data quality, filtering, and discretization")

    # Left column - DQ
    tf = _add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.2), "", font_size=13)
    _section_heading(tf, "Data Quality & Filtering", font_size=16)
    _bullet(tf, "Schema validation, missing values, date range checks")
    _bullet(tf, "Fraud / out-of-norm / legal entity filtering")
    _bullet(tf, "Outlier flagging via Z-score (configurable threshold)")
    _bullet(tf, "Status / reject-reason reclassification")
    _bullet(tf, "Monthly distribution logging")

    # Right column - Direction
    tf2 = _add_textbox(slide, Inches(7), Inches(2.0), Inches(6), Inches(2.2), "", font_size=13)
    _section_heading(tf2, "Monotonicity Direction Inference", font_size=16)
    _add_para(tf2, "Weighted Spearman correlation on booked data (weighted by production).",
              font_size=12, space_before=Pt(6))
    _bullet(tf2, "Direction = -1: higher bin = safer (e.g., credit scores)", font_size=12)
    _bullet(tf2, "Direction = +1: higher bin = riskier", font_size=12)
    _add_para(tf2, "Can be overridden via [preprocessing.directions] config.",
              font_size=11, color=GREY, space_before=Pt(6))

    # Binning table
    tf3 = _add_textbox(slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(0.5), "", font_size=13)
    _section_heading(tf3, "Binning Methods", font_size=16)

    _add_table(slide, Inches(0.8), Inches(4.8), Inches(11.5), [
        ["Method", "Algorithm", "Best For"],
        ["quantile (default)", "Equal-count quantile splits", "Unsupervised, avoids target leakage"],
        ["optimization", "DecisionTreeRegressor + production weights", "Maximize risk separation across bins"],
    ], font_size=12)

    # BinConfig
    tf4 = _add_textbox(slide, Inches(0.8), Inches(5.8), Inches(5.5), Inches(0.5), "", font_size=13)
    _section_heading(tf4, "BinConfig (per variable)", font_size=16)

    _code_box(slide, Inches(0.8), Inches(6.3), Inches(6), Inches(0.65),
              "source_col: str     # Raw column (e.g., 'income_t1t2_m')\n"
              "output_col: str     # Binned column (e.g., 'income_bin')\n"
              "bin_edges | max_bins # Fixed edges or supervised learning", font_size=11)

    # Gini assessment
    tf5 = _add_textbox(slide, Inches(7), Inches(5.8), Inches(6), Inches(1.5), "", font_size=13)
    _section_heading(tf5, "Gini Retention Assessment", font_size=16)
    _add_para(tf5, "Compares AUC-based Gini for raw score vs. binned score.",
              font_size=12, space_before=Pt(6))
    _add_para(tf5, "Reports gini_raw, gini_binned, and retention %.",
              font_size=12, space_before=Pt(3))

    _draw_footer(slide, 5, TOTAL_SLIDES)


def slide_risk_adjustment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_header(slide, "Risk Adjustment Pipeline", "Three sequential corrections for the rejected population")

    # Flow diagram
    labels = ["Model\nPrediction", "Stress\nFactor", "Reject\nInference", "Tasa Fin\n(Transf. Rate)"]
    for i, label in enumerate(labels):
        x = Inches(0.8 + i * 3.1)
        color = [LIGHT_BLUE, MID_BLUE, DARK_BLUE, RGBColor(0x27, 0xAE, 0x60)][i]
        _flow_box(slide, x, Inches(2.1), Inches(2.5), Inches(0.7), label, bg_color=color)
        if i < 3:
            _arrow(slide, x + Inches(2.6), Inches(2.25))

    # Stress factor section
    tf = _add_textbox(slide, Inches(0.8), Inches(3.2), Inches(5.8), Inches(3.5), "", font_size=13)
    _section_heading(tf, "Stress Factor (3 modes)", font_size=16)
    _bullet(tf, "stress_factor = worst_5%_risk / overall_risk", bold_prefix="Global:", font_size=12)
    _bullet(tf, "Separate stress per grid cell (fallback to global for bins < 20 obs)", bold_prefix="Per-bin:", font_size=12)
    _bullet(tf, "No stress (factor = 1.0). Recommended when parceling is active", bold_prefix="Disabled:", font_size=12)
    _add_para(tf, "", font_size=6, space_before=Pt(2))
    _section_heading(tf, "Transformation Rate (tasa_fin)", font_size=16)
    _bullet(tf, "tasa_fin = booked_amount / eligible_amount", font_size=12)
    _bullet(tf, "Applied to ALL repesca indicators after reject inference", font_size=12)
    _bullet(tf, "Per-bin mode: computes tasa_fin per grid cell (bins < 10 obs fall back to global)", bold_prefix="per_bin_tasa_fin:", font_size=12)

    # Integration sequence
    tf2 = _add_textbox(slide, Inches(7), Inches(3.2), Inches(6), Inches(3.8), "", font_size=13)
    _section_heading(tf2, "Integration Sequence", font_size=16)
    steps = [
        "1. Aggregate repesca by N-D grid (sum indicators x annual_coef)",
        "2. Predict risk via trained model (stress embedded)",
        "3. Apply reject inference parceling multiplier",
        "4. Scale ALL indicators by tasa_fin (global or per-bin)",
        "5. Merge with booked summary \u2192 CellGrid \u2192 MILP",
    ]
    for step in steps:
        _bullet(tf2, step, font_size=12)

    _add_para(tf2, "", font_size=4, space_before=Pt(4))
    _add_para(tf2, "Booked data: no adjustments needed -- observed values used directly.",
              font_size=11, color=GREY, space_before=Pt(6))
    _add_para(tf2, "Double-counting warning: using global stress + parceling may overcorrect. "
              "Consider stress_mode='disabled' when parceling is active.",
              font_size=11, color=RGBColor(0xC0, 0x39, 0x2B), space_before=Pt(4))

    _draw_footer(slide, 6, TOTAL_SLIDES)


def slide_inference(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_header(slide, "Model Inference", "Custom estimators, feature engineering, cross-validation")

    # Estimators
    tf = _add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(2.5), "", font_size=13)
    _section_heading(tf, "Estimators Evaluated", font_size=16)

    _add_table(slide, Inches(0.8), Inches(2.5), Inches(5.5), [
        ["Estimator", "Key Property"],
        ["LinearRegression", "Baseline, fully interpretable"],
        ["Ridge / Lasso / ElasticNet", "Regularized linear models"],
        ["HurdleRegressor", "Two-stage: P(Y>0) x E[Y|Y>0]"],
        ["TweedieGLM", "Compound Poisson-Gamma (power=1.5)"],
        ["XGBoost", "Monotonic constraints, Optuna tuning"],
        ["LightGBM", "Monotonic constraints, Optuna tuning"],
    ], font_size=10, header_font_size=11)

    # Feature engineering
    tf2 = _add_textbox(slide, Inches(7), Inches(2.0), Inches(6), Inches(0.5), "", font_size=13)
    _section_heading(tf2, "Feature Engineering & Selection", font_size=16)

    _add_table(slide, Inches(7), Inches(2.5), Inches(5.8), [
        ["Feature Set", "Description"],
        ["degree_1", "Raw variables only"],
        ["degree_2", "Interactions + squared terms"],
        ["degree_3", "Full cubic polynomial + interactions"],
    ], font_size=11)

    # CV strategy
    tf3 = _add_textbox(slide, Inches(0.8), Inches(5.2), Inches(5.8), Inches(2.2), "", font_size=13)
    _section_heading(tf3, "Cross-Validation Strategy", font_size=16)
    _bullet(tf3, "K-Fold CV (configurable, default k=4)", font_size=12)
    _bullet(tf3, "Validation folds split on raw data BEFORE aggregation", font_size=12)
    _bullet(tf3, "RMSE weighted by group size (exposure-weighted)", font_size=12)
    _bullet(tf3, "Optuna HPO for tree models (XGBoost, LightGBM)", font_size=12)

    # 1SE rule
    tf4 = _add_textbox(slide, Inches(7), Inches(4.5), Inches(6), Inches(2.5), "", font_size=13)
    _section_heading(tf4, "One-Standard-Error Rule", font_size=16)
    _add_para(tf4, "Among models within 1 SE of the best CV RMSE, "
              "select the simplest. Applied twice:",
              font_size=12, space_before=Pt(6))
    _bullet(tf4, "Model type selection (Linear vs Ridge vs XGBoost...)", font_size=12)
    _bullet(tf4, "Feature set selection (degree_1 vs degree_2 vs degree_3)", font_size=12)
    _add_para(tf4, "Prevents overfitting from unnecessarily complex models.",
              font_size=11, color=GREY, space_before=Pt(6))

    _draw_footer(slide, 7, TOTAL_SLIDES)


def slide_optimization(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_header(slide, "MILP Optimization", "Maximize production subject to risk budget and monotonicity")

    # MILP formulation
    _code_box(slide, Inches(0.8), Inches(2.1), Inches(7.5), Inches(2.0),
              "Decision:     x[i] in {0, 1}   for each cell i in the N-D grid\n\n"
              "Objective:    max  SUM( oa_amt_h0[i] * x[i] )\n\n"
              "Risk budget:  SUM( (mult * todu_30ever_h6[i]\n"
              "                   - target_risk/100 * todu_amt_pile_h6[i]) * x[i] ) <= 0\n\n"
              "Monotonicity: x[riskier] - x[safer] <= 0   (per dimension, per adjacent pair)",
              font_size=12)

    # Monotonicity
    tf = _add_textbox(slide, Inches(9), Inches(2.1), Inches(4), Inches(2.0), "", font_size=13)
    _section_heading(tf, "Monotonicity", font_size=16)
    _add_para(tf, 'If a riskier cell is rejected, all cells riskier along '
              'that dimension must also be rejected.',
              font_size=12, space_before=Pt(6))
    _add_para(tf, 'Creates a "staircase" acceptance pattern. Enforced marginally '
              '(per dimension independently).',
              font_size=11, color=GREY, space_before=Pt(6))

    # Swap-in constraints
    tf2 = _add_textbox(slide, Inches(0.8), Inches(4.4), Inches(5.8), Inches(2.5), "", font_size=13)
    _section_heading(tf2, "Optional MILP Constraints", font_size=16)
    _bullet(tf2, "Cap % of production from swap-in (repesca)", bold_prefix="Swap-in production:", font_size=12)
    _bullet(tf2, "Cap risk of swap-in population alone", bold_prefix="Swap-in risk:", font_size=12)
    _bullet(tf2, "Pin cells as forced-accept or forced-reject", bold_prefix="Fixed cells:", font_size=12)

    # Solver
    tf3 = _add_textbox(slide, Inches(7), Inches(4.4), Inches(6), Inches(2.5), "", font_size=13)
    _section_heading(tf3, "Solver & Fallback", font_size=16)
    _bullet(tf3, "scipy.optimize.milp (configurable timeout, default 30s)", bold_prefix="Primary:", font_size=12)
    _bullet(tf3, "pymoo Genetic Algorithm (N > 2 variables)", bold_prefix="Fallback:", font_size=12)
    _bullet(tf3, "Legacy brute-force enumeration (2 variables)", bold_prefix="Legacy:", font_size=12)
    _add_para(tf3, "Built as sparse CSC matrix for ~100x faster constraint setup.",
              font_size=11, color=GREY, space_before=Pt(8))

    _draw_footer(slide, 8, TOTAL_SLIDES)


def slide_pareto(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_header(slide, "Pareto Frontier", "Risk vs. production trade-off analysis")

    tf = _add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(3.0), "", font_size=13)
    _section_heading(tf, "Frontier Construction", font_size=16)
    _bullet(tf, "50 linearly-spaced risk targets from 0.01% to max_risk x 1.1", bold_prefix="Step 1:", font_size=12)
    _bullet(tf, "MILP solve at each target \u2192 (risk, production, mask) tuples", bold_prefix="Step 2:", font_size=12)
    _bullet(tf, "Deduplicate by unique binary acceptance masks", bold_prefix="Step 3:", font_size=12)
    _bullet(tf, "Sort by ascending risk, filter to strictly increasing production", bold_prefix="Step 4:", font_size=12)
    _bullet(tf, "Full dominance verification: remove any remaining dominated points", bold_prefix="Step 5:", font_size=12)

    # Scenario table
    tf2 = _add_textbox(slide, Inches(7), Inches(2.0), Inches(6), Inches(0.5), "", font_size=13)
    _section_heading(tf2, "Scenario Selection", font_size=16)

    _add_table(slide, Inches(7), Inches(2.5), Inches(5.8), [
        ["Scenario", "Risk Target", "Interpretation"],
        ["Pessimistic", "optimum_risk - risk_step", "Conservative: lower risk tolerance"],
        ["Base", "optimum_risk", "Central estimate: target risk level"],
        ["Optimistic", "optimum_risk + risk_step", "Aggressive: higher risk tolerance"],
    ], font_size=11)

    # Output
    tf3 = _add_textbox(slide, Inches(0.8), Inches(5.3), Inches(12), Inches(1.8), "", font_size=13)
    _section_heading(tf3, "Per-Scenario Outputs", font_size=16)
    _bullet(tf3, "Risk/production summary table (Actual vs Optimum vs Swap-in/Swap-out)", font_size=12)
    _bullet(tf3, "Optimal solution with per-variable cutoff values and KPIs", font_size=12)
    _bullet(tf3, "Interactive Pareto dashboard (Plotly HTML) with selected point highlighted", font_size=12)
    _bullet(tf3, "Acceptance grid heatmap and cutoff summary (wide + long format)", font_size=12)

    _draw_footer(slide, 9, TOTAL_SLIDES)


def slide_reject_inference(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_header(slide, "Reject Inference", "Correcting selection bias in observed outcomes")

    # Problem
    tf = _add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(1.5), "", font_size=13)
    _section_heading(tf, "The Problem", font_size=16)
    _add_para(tf, "Risk models trained on booked data systematically underestimate "
              "risk for rejected applicants. Reject inference corrects this bias "
              "proportionally to how selective each bin's original acceptance was.",
              font_size=12, space_before=Pt(6))

    # Methods table
    tf2 = _add_textbox(slide, Inches(0.8), Inches(3.6), Inches(5.8), Inches(0.5), "", font_size=13)
    _section_heading(tf2, "Parceling Methods", font_size=16)

    _add_table(slide, Inches(0.8), Inches(4.1), Inches(11.5), [
        ["Method", "Formula", "Behavior"],
        ["Linear (default)", "mult = 1 + factor x (1 - AR)", "Linear increase as acceptance rate drops"],
        ["Power", "mult = (1 / AR) ^ factor", "Exponential growth at low acceptance rates"],
        ["Sigmoid", "mult = 1 + factor / (1 + exp(10 x (AR - 0.5)))", "Smooth S-curve, steep around 50% AR"],
    ], font_size=11)

    # Enhancements
    tf3 = _add_textbox(slide, Inches(7), Inches(2.0), Inches(6), Inches(1.5), "", font_size=13)
    _section_heading(tf3, "Enhancements", font_size=16)
    _bullet(tf3, "Beta-Binomial smoothing of noisy acceptance rates", bold_prefix="Bayesian smoothing:", font_size=12)
    _bullet(tf3, "Isotonic regression per variable axis", bold_prefix="Monotonicity:", font_size=12)
    _bullet(tf3, "1 - exp(-n/50) per bin", bold_prefix="Confidence scores:", font_size=12)

    # Algorithm
    tf4 = _add_textbox(slide, Inches(0.8), Inches(5.5), Inches(12), Inches(1.5), "", font_size=13)
    _section_heading(tf4, "Algorithm", font_size=16)
    _add_para(tf4,
              "AR[bin] = n_booked / (n_booked + n_score_rejected)  \u2192  "
              "multiplier = f(AR)  \u2192  "
              "cap at max_risk_multiplier  \u2192  "
              "todu_30ever_h6 *= multiplier",
              font_size=12, space_before=Pt(6))
    _add_para(tf4,
              "Only risk is adjusted. Production (oa_amt) is observable for all applications. "
              "Multipliers floored at 1.0 (risk never adjusted downward). "
              "Parameters auto-tunable via RI optimizer (grid search or Optuna).",
              font_size=11, color=GREY, space_before=Pt(4))

    _draw_footer(slide, 10, TOTAL_SLIDES)


def slide_scenarios_bootstrap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_header(slide, "Scenario Analysis & Bootstrap CI", "Quantifying uncertainty and validating cutoffs")

    # Bootstrap
    tf = _add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(2.5), "", font_size=13)
    _section_heading(tf, "Bootstrap Confidence Intervals", font_size=16)
    _add_para(tf, "Record-level resampling of the booked portfolio (default: 1,000 iterations). "
              "At each resample, recalculates production and risk respecting the selected cutoff.",
              font_size=12, space_before=Pt(6))
    _bullet(tf, "2.5th and 97.5th percentiles \u2192 95% CI", font_size=12)
    _bullet(tf, "Standard error reported alongside point estimates", font_size=12)
    _bullet(tf, "Configurable: n_bootstraps (100 - 50,000)", font_size=12)

    # Audit tables
    tf2 = _add_textbox(slide, Inches(7), Inches(2.0), Inches(6), Inches(2.5), "", font_size=13)
    _section_heading(tf2, "Audit Tables (Record-Level)", font_size=16)

    _add_table(slide, Inches(7), Inches(2.5), Inches(5.8), [
        ["Classification", "Description"],
        ["keep", "Booked and passes proposed cutoff"],
        ["swap_out", "Booked but fails proposed cutoff"],
        ["swap_in", "Score-rejected but passes cutoff"],
        ["rejected", "Score-rejected and still fails"],
        ["rejected_other", "Policy rejections (not candidates)"],
    ], font_size=10, header_font_size=11)

    # RI optimizer
    tf3 = _add_textbox(slide, Inches(0.8), Inches(4.8), Inches(12), Inches(2.2), "", font_size=13)
    _section_heading(tf3, "RI Parameter Optimizer (Optional)", font_size=16)
    _add_para(tf3, "Automated search over (uplift_factor, max_risk_multiplier) to minimize "
              "calibration error against the selection-bias model: target = booked_risk / AR^gamma.",
              font_size=12, space_before=Pt(6))
    _bullet(tf3, "Grid search (exhaustive, default 11 x 9 = 99 combos) or Optuna TPE", font_size=12)
    _bullet(tf3, "Selection: min calibration error among feasible solutions; ties within 5% broken by max production", font_size=12)
    _bullet(tf3, "MR validation: reports degradation_ratio = mr_error / main_error (near 1.0 = temporally stable)", font_size=12)
    _bullet(tf3, "If better params found, re-runs full optimization pipeline automatically (Step 6d)", font_size=12)

    _draw_footer(slide, 11, TOTAL_SLIDES)


def slide_stability_mr(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_header(slide, "Stability & MR Validation", "PSI/CSI monitoring and out-of-time validation")

    # PSI
    tf = _add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(0.5), "", font_size=13)
    _section_heading(tf, "Population Stability Index (PSI)", font_size=16)

    _code_box(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(0.45),
              "PSI = SUM( (Actual% - Expected%) x ln(Actual% / Expected%) )",
              font_size=12)

    _add_table(slide, Inches(0.8), Inches(3.1), Inches(5.5), [
        ["PSI Range", "Interpretation", "Action"],
        ["< 0.10", "Stable", "No action needed"],
        ["0.10 - 0.25", "Moderate shift", "Monitor closely"],
        [">= 0.25", "Significant shift", "Investigate & recalibrate"],
    ], font_size=11)

    # MR Validation
    tf2 = _add_textbox(slide, Inches(7), Inches(2.0), Inches(6), Inches(2.5), "", font_size=13)
    _section_heading(tf2, "MR (Out-of-Time) Validation", font_size=16)
    _bullet(tf2, "Apply optimized cutoffs to recent holdout period", font_size=12)
    _bullet(tf2, "Compare realized risk/production vs. main period estimates", font_size=12)
    _bullet(tf2, "Swap-in / swap-out breakdown per scenario", font_size=12)
    _bullet(tf2, "PSI computed per variable between main and MR periods", font_size=12)

    # Hybrid risk
    tf3 = _add_textbox(slide, Inches(0.8), Inches(4.7), Inches(12), Inches(2.5), "", font_size=13)
    _section_heading(tf3, "Hybrid MR Risk Estimation (4-tier fallback)", font_size=16)

    _add_table(slide, Inches(0.8), Inches(5.2), Inches(11.5), [
        ["Priority", "Source", "Condition"],
        ["1", "MR observed H6", "Enough valid MR H6 observations in the bin"],
        ["2", "H3 extrapolated", "H3 mature, valid main-period H6/H3 ratio available"],
        ["3", "Main imputed", "Main-period bin exists but MR evidence insufficient"],
        ["4", "Model fallback", "Bin only exists in MR, no observed outcomes"],
    ], font_size=11)

    _add_para(tf3, "", font_size=2, space_before=Pt(1))
    _add_para(tf3, "H3\u2192H6 extrapolation methods: linear, power, logistic, or auto "
              "(fits curvature via weighted log-log regression from main-period data).",
              font_size=11, color=GREY, space_before=Pt(60))

    _draw_footer(slide, 12, TOTAL_SLIDES)


def slide_sensitivity(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_header(slide, "Sensitivity, Marginal Impact & Fixed Cells",
                 "Robustness analysis and interactive re-optimization")

    # Sensitivity
    tf = _add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(2.5), "", font_size=13)
    _section_heading(tf, "Sensitivity Analysis", font_size=16)
    _add_para(tf, "Perturbs risk indicator (todu_30ever_h6_rep) by configurable "
              "percentages and re-solves the MILP at each level.",
              font_size=12, space_before=Pt(6))
    _bullet(tf, "Default levels: +/-5%, 10%, 20%, 30%", font_size=12)
    _bullet(tf, "Counts cells that flip accept\u2192reject and reject\u2192accept", font_size=12)
    _bullet(tf, "Per-cell flip threshold: minimum perturbation to change status", font_size=12)

    # Marginal impact
    tf2 = _add_textbox(slide, Inches(7), Inches(2.0), Inches(6), Inches(2.5), "", font_size=13)
    _section_heading(tf2, "Marginal Impact (O(N) Analytical)", font_size=16)
    _add_para(tf2, "For each cell, computes the exact production and risk change "
              "from flipping its accept/reject status.",
              font_size=12, space_before=Pt(6))
    _bullet(tf2, "delta_production: absolute change in portfolio production", font_size=12)
    _bullet(tf2, "delta_risk_pct: percentage point change in b2_ever_h6", font_size=12)
    _add_para(tf2, "Computed analytically from baseline sums -- no re-solving needed.",
              font_size=11, color=GREY, space_before=Pt(6))

    # Fixed-cell
    tf3 = _add_textbox(slide, Inches(0.8), Inches(4.8), Inches(12), Inches(2.2), "", font_size=13)
    _section_heading(tf3, "Fixed-Cell Constraints (Pin Mode)", font_size=16)
    _add_para(tf3, "Individual cells can be pinned as forced-accept or forced-reject. "
              "The MILP enforces pins by setting lb = ub on those variables. "
              "Monotonicity and risk constraints remain active.",
              font_size=12, space_before=Pt(6))
    _bullet(tf3, "Programmatic: solve_with_fixed_cells() in optimization_utils.py", font_size=12)
    _bullet(tf3, "Interactive: Pin Mode in dashboard Cutoff Explorer (click cells to cycle: unpinned \u2192 accept \u2192 reject \u2192 unpinned)", font_size=12)
    _add_para(tf3, "Returns None if pins are contradictory or make risk target infeasible.",
              font_size=11, color=GREY, space_before=Pt(6))

    _draw_footer(slide, 13, TOTAL_SLIDES)


def slide_batch_allocation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_header(slide, "Batch Processing & Global Allocation",
                 "Multi-segment execution and portfolio-level optimization")

    # Batch
    tf = _add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(2.8), "", font_size=13)
    _section_heading(tf, "Batch Runner (run_batch.py)", font_size=16)
    _bullet(tf, "Two-tier config: config.toml + per-segment segments.toml overrides", font_size=12)
    _bullet(tf, "Global bin learning: learn edges once on full population, share across segments", font_size=12)
    _bullet(tf, "Parallel execution via ProcessPoolExecutor", font_size=12)
    _bullet(tf, "Supersegment training: train on combined data, optimize per segment", font_size=12)
    _bullet(tf, "Consolidated HTML report with cross-segment comparison", font_size=12)

    _code_box(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(0.9),
              "uv run python run_batch.py                    # all segments\n"
              "uv run python run_batch.py -s seg1 seg2       # specific\n"
              "uv run python run_batch.py --parallel -w 4    # parallel",
              font_size=11)

    # Allocation
    tf2 = _add_textbox(slide, Inches(7), Inches(2.0), Inches(6), Inches(3.0), "", font_size=13)
    _section_heading(tf2, "Global Allocation (run_allocation.py)", font_size=16)
    _add_para(tf2, "Given N segment Pareto frontiers, select one operating point per "
              "segment to maximize global production subject to a weighted-average "
              "global risk constraint.",
              font_size=12, space_before=Pt(6))

    _add_table(slide, Inches(7), Inches(3.3), Inches(5.8), [
        ["Method", "Algorithm", "Trade-off"],
        ["Exact (MILP)", "Binary vars per (segment, solution)", "Globally optimal"],
        ["Greedy", "Hill-climbing heuristic", "Faster, local optima"],
    ], font_size=11)

    # Constraints
    tf3 = _add_textbox(slide, Inches(7), Inches(4.5), Inches(6), Inches(2.5), "", font_size=13)
    _section_heading(tf3, "Per-Segment Constraints", font_size=16)
    _bullet(tf3, "min_risk, max_risk bounds", font_size=12)
    _bullet(tf3, "min_production (production floor)", font_size=12)
    _bullet(tf3, "locked_sol_fac (lock to specific frontier point)", font_size=12)
    _bullet(tf3, "Portfolio-wide --production-floor flag", font_size=12)

    _draw_footer(slide, 14, TOTAL_SLIDES)


def slide_trends(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_header(slide, "Trend Analysis & SPC Anomaly Detection",
                 "Monthly monitoring with statistical process control")

    # Metrics
    tf = _add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(2.5), "", font_size=13)
    _section_heading(tf, "Monthly Metric Aggregation", font_size=16)
    _bullet(tf, "Approval rate, total records, booked count", font_size=12)
    _bullet(tf, "Mean production, total production", font_size=12)
    _bullet(tf, "Risk rate per score bin", font_size=12)
    _bullet(tf, "Score distribution means", font_size=12)
    _add_para(tf, "Output: monthly_metrics_{segment}.csv + interactive Plotly charts.",
              font_size=11, color=GREY, space_before=Pt(8))

    # SPC
    tf2 = _add_textbox(slide, Inches(7), Inches(2.0), Inches(6), Inches(3.0), "", font_size=13)
    _section_heading(tf2, "Robust SPC Anomaly Detection", font_size=16)
    _add_para(tf2, "Uses a one-period-lagged rolling estimator to avoid lookahead bias:",
              font_size=12, space_before=Pt(6))

    _code_box(slide, Inches(7), Inches(3.0), Inches(5.8), Inches(1.2),
              "Center:  median(x_{t-w}, ..., x_{t-1})     [lagged]\n"
              "Scale:   MAD x 1.4826  (std-equivalent)    [robust]\n"
              "Threshold: t-distribution critical value    [small-sample]\n"
              "Anomaly:  |x_t - center| > threshold x scale",
              font_size=11)

    _add_para(tf2, "", font_size=2, space_before=Pt(1))
    _add_para(tf2, "Robust to outliers (median + MAD instead of mean + std). "
              "Student-t critical value for short windows instead of fixed z-score.",
              font_size=11, color=GREY, space_before=Pt(56))

    # Score discriminance
    tf3 = _add_textbox(slide, Inches(0.8), Inches(5.0), Inches(12), Inches(2.0), "", font_size=13)
    _section_heading(tf3, "Score Discriminance (run_score_metrics.py)", font_size=16)
    _bullet(tf3, "Gini coefficient and KS statistic per score variable", font_size=12)
    _bullet(tf3, "Lift tables (decile-based), precision-recall and ROC curves", font_size=12)
    _bullet(tf3, "DeLong test for pairwise statistical comparison between score models", font_size=12)
    _bullet(tf3, "Per-segment, per-supersegment, main period vs. MR period comparison", font_size=12)

    _draw_footer(slide, 15, TOTAL_SLIDES)


def slide_reporting(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_header(slide, "Consolidated Reporting", "Segment reports and cross-segment comparison")

    # Per-segment report
    tf = _add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(3.0), "", font_size=13)
    _section_heading(tf, "Per-Segment HTML Report", font_size=16)
    _add_para(tf, "Self-contained Jinja2 template with embedded Plotly charts:",
              font_size=12, space_before=Pt(6))
    _bullet(tf, "Risk vs. production scatter (preprocessing)", font_size=12)
    _bullet(tf, "Pareto frontier with scenario selection", font_size=12)
    _bullet(tf, "Acceptance grid heatmaps", font_size=12)
    _bullet(tf, "MR validation comparison", font_size=12)
    _bullet(tf, "PSI/CSI stability dashboard", font_size=12)
    _bullet(tf, "Monthly trend charts with anomaly markers", font_size=12)

    # Consolidated
    tf2 = _add_textbox(slide, Inches(7), Inches(2.0), Inches(6), Inches(3.0), "", font_size=13)
    _section_heading(tf2, "Consolidated Cross-Segment Report", font_size=16)
    _add_para(tf2, "Generated after batch processing:", font_size=12, space_before=Pt(6))
    _bullet(tf2, "Scenario KPI summary: 1 row per (segment x scenario)", bold_prefix="Table:", font_size=12)
    _bullet(tf2, "Colored acceptance matrices per segment (green/red/grey)", bold_prefix="Grids:", font_size=12)
    _bullet(tf2, "N-variable slicing with sub-tables (max 12 combinations)", bold_prefix="N-D:", font_size=12)
    _bullet(tf2, "Portfolio-level aggregated metrics", bold_prefix="Portfolio:", font_size=12)

    # Dashboards
    tf3 = _add_textbox(slide, Inches(0.8), Inches(5.2), Inches(12), Inches(2.0), "", font_size=13)
    _section_heading(tf3, "Interactive Dashboards", font_size=16)

    _add_table(slide, Inches(0.8), Inches(5.7), Inches(11.5), [
        ["Dashboard", "Entry Point", "Features"],
        ["Results Explorer", "dashboard.py (Dash, port 8050)", "Scenario comparison, cutoff explorer, pin mode, uncertainty overlays"],
        ["Global Allocator", "interactive_allocator.py (Dash, port 8051)", "Real-time portfolio allocation, per-segment risk sliders"],
        ["Gradio UI", "gradio_dashboard.py", "Simplified results exploration"],
    ], font_size=10, header_font_size=11)

    _draw_footer(slide, 16, TOTAL_SLIDES)


def slide_config_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _draw_header(slide, "Configuration System", "Two-tier Pydantic-validated config with 50+ parameters")

    tf = _add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(0.5), "", font_size=13)
    _section_heading(tf, "config.toml -- Key Parameter Groups", font_size=16)

    _add_table(slide, Inches(0.8), Inches(2.5), Inches(5.5), [
        ["Group", "Key Parameters"],
        ["Core", "variables, indicators, keep_vars, dates"],
        ["Binning", "bins.*, octroi_bins/efx_bins, directions"],
        ["Economic", "multiplier, optimum_risk, risk_step, n_months"],
        ["Stress", "stress_mode (global / per_bin / disabled)"],
        ["Tasa Fin", "per_bin_tasa_fin (global / per-bin)"],
        ["Reject Inf.", "method, parceling, uplift, max_mult, Bayesian"],
        ["RI Optimizer", "method (grid/optuna), gamma, ranges, steps"],
        ["MR Period", "dates, use_mr_outcomes, extrapolation_method"],
        ["MILP Tuning", "milp_time_limit, pareto_n_points, n_bootstraps"],
        ["Constraints", "max_swapin_production_pct, max_swapin_risk"],
        ["Sensitivity", "run_sensitivity, sensitivity_levels"],
    ], font_size=10, header_font_size=11)

    # segments.toml
    tf2 = _add_textbox(slide, Inches(7), Inches(2.0), Inches(6), Inches(0.5), "", font_size=13)
    _section_heading(tf2, "segments.toml -- Per-Segment Overrides", font_size=16)
    _add_para(tf2, "", font_size=2, space_before=Pt(1))
    _add_para(tf2, "Any config.toml field can be overridden per segment via recursive merge.",
              font_size=12, space_before=Pt(6))

    _add_table(slide, Inches(7), Inches(3.0), Inches(5.8), [
        ["Key", "Purpose"],
        ["segment_filter", "Segment regex (required)"],
        ["optimum_risk / risk_step", "Per-segment risk appetite"],
        ["supersegment", "Shared model name"],
        ["min_risk / max_risk", "Allocation bounds"],
        ["min_production", "Production floor"],
        ["locked_sol_fac", "Lock to frontier point"],
        ["fixed_cutoffs", "Skip MILP, use predefined cutoffs"],
        ["bins.*, directions", "Override binning per segment"],
    ], font_size=10, header_font_size=11)

    # Validation
    tf3 = _add_textbox(slide, Inches(7), Inches(5.7), Inches(6), Inches(1.5), "", font_size=13)
    _section_heading(tf3, "Pydantic Validation", font_size=16)
    _bullet(tf3, "Field validators: non-empty lists, date formats, ranges", font_size=11)
    _bullet(tf3, "Model validators: date ordering, MR pair consistency, bin auto-population", font_size=11)
    _bullet(tf3, "variables >= 2, unique, inference_variables subset of variables", font_size=11)

    _draw_footer(slide, 17, TOTAL_SLIDES)


def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BLUE)

    _add_textbox(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
                 "Summary", font_size=40, color=WHITE, bold=True)

    _add_textbox(slide, Inches(1), Inches(1.6), Inches(11), Inches(0.5),
                 "A complete, production-ready credit decisioning pipeline",
                 font_size=18, color=RGBColor(0xBB, 0xCC, 0xDD))

    bar = slide.shapes.add_shape(1, Inches(0), Inches(2.3), SLIDE_WIDTH, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()

    items = [
        "Multi-segment batch processing with parallel execution and supersegment models",
        "N-dimensional score grids with quantile and optimization-aware binning",
        "7 ML estimators (Hurdle, Tweedie, XGBoost, LightGBM, Ridge, Lasso, ElasticNet) with CV + 1SE rule",
        "MILP optimization with monotonicity, swap-in, and fixed-cell constraints on Pareto frontier",
        "Reject inference (3 methods) with Bayesian smoothing, monotonicity, and automated parameter tuning",
        "Configurable stress factor (global / per-bin / disabled) and per-bin transformation rate",
        "H3\u2192H6 extrapolation with auto-calibrated curvature for MR validation",
        "Bootstrap CI, PSI/CSI stability, sensitivity analysis, and marginal impact",
        "Robust SPC trend monitoring with median + MAD + Student-t anomaly detection",
        f"799 tests, Pydantic config, 50+ configurable parameters, Docker-ready deployment",
    ]

    tf = _add_textbox(slide, Inches(1.5), Inches(2.7), Inches(10.5), Inches(4.5), "", font_size=15)
    for item in items:
        _add_para(tf, f"  {item}", font_size=15, color=WHITE, space_before=Pt(10))

    _draw_footer(slide, 18, TOTAL_SLIDES)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_title(prs)
    slide_agenda(prs)
    slide_architecture(prs)
    slide_populations_risk(prs)
    slide_preprocessing(prs)
    slide_risk_adjustment(prs)
    slide_inference(prs)
    slide_optimization(prs)
    slide_pareto(prs)
    slide_reject_inference(prs)
    slide_scenarios_bootstrap(prs)
    slide_stability_mr(prs)
    slide_sensitivity(prs)
    slide_batch_allocation(prs)
    slide_trends(prs)
    slide_reporting(prs)
    slide_config_summary(prs)
    slide_summary(prs)

    out = Path("Credit_Risk_Scoring_Pipeline_Presentation.pptx")
    prs.save(str(out))
    print(f"Saved: {out}")
    return out


def convert_to_pdf(pptx_path: Path) -> Path:
    pdf_path = pptx_path.with_suffix(".pdf")
    try:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", str(pptx_path)],
            check=True, capture_output=True, timeout=60,
        )
        print(f"Saved: {pdf_path}")
        return pdf_path
    except FileNotFoundError:
        print("LibreOffice not found. Install it to auto-convert to PDF.")
        print(f"Manual conversion: open {pptx_path} and export as PDF.")
        sys.exit(1)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate pipeline presentation")
    parser.add_argument("--pdf", action="store_true", help="Also convert to PDF (requires LibreOffice)")
    args = parser.parse_args()

    pptx_path = build_presentation()
    if args.pdf:
        convert_to_pdf(pptx_path)
